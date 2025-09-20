import os
import tempfile
import shutil
from pptx import Presentation
from combine_ppt import combine_ppt

def create_sample_pptx(file_path, slide_texts):
    prs = Presentation()
    for i, text in enumerate(slide_texts):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = text
    prs.save(file_path)

def test_combine_ppt_basic():
    with tempfile.TemporaryDirectory() as tmpdir:
        ppt1 = os.path.join(tmpdir, "ppt1.pptx")
        ppt2 = os.path.join(tmpdir, "ppt2.pptx")
        output = os.path.join(tmpdir, "combined.pptx")

        create_sample_pptx(ppt1, ["Slide 1A", "Slide 1B"])
        create_sample_pptx(ppt2, ["Slide 2A"])

        combine_ppt([ppt1, ppt2], output)

        combined = Presentation(output)
        slide_texts = [slide.shapes.title.text for slide in combined.slides]
        assert slide_texts == ["Slide 1A", "Slide 1B", "Slide 2A"]

def test_combine_ppt_empty_input():
    with tempfile.TemporaryDirectory() as tmpdir:
        output = os.path.join(tmpdir, "combined_empty.pptx")
        combine_ppt([], output)
        prs = Presentation(output)
        # Should have only the default slide if any, or none
        assert len(prs.slides) == 0

def test_combine_ppt_preserves_order():
    with tempfile.TemporaryDirectory() as tmpdir:
        ppt1 = os.path.join(tmpdir, "ppt1.pptx")
        ppt2 = os.path.join(tmpdir, "ppt2.pptx")
        ppt3 = os.path.join(tmpdir, "ppt3.pptx")
        output = os.path.join(tmpdir, "combined_order.pptx")

        create_sample_pptx(ppt1, ["A"])
        create_sample_pptx(ppt2, ["B"])
        create_sample_pptx(ppt3, ["C"])

        combine_ppt([ppt2, ppt1, ppt3], output)
        combined = Presentation(output)
        slide_texts = [slide.shapes.title.text for slide in combined.slides]
        assert slide_texts == ["B", "A", "C"]

def test_combine_ppt_handles_nonexistent_file():
    with tempfile.TemporaryDirectory() as tmpdir:
        ppt1 = os.path.join(tmpdir, "ppt1.pptx")
        output = os.path.join(tmpdir, "combined.pptx")
        create_sample_pptx(ppt1, ["Slide"])
        try:
            combine_ppt([ppt1, "nonexistent.pptx"], output)
            assert False, "Should raise an exception for nonexistent file"
        except Exception:
            assert True