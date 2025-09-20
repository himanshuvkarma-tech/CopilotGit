def combine_ppt(ppt_files, output_file):
    """
    Combines multiple PowerPoint (.pptx) files into a single presentation.

    This function takes a list of PowerPoint file paths, merges all slides from each file into a new presentation,
    and saves the combined presentation to the specified output file.

    Args:
        ppt_files (list of str): List of file paths to the PowerPoint presentations to combine.
        output_file (str): File path where the combined presentation will be saved.

    Notes:
        - Slide layouts and placeholders are preserved where possible.
        - Non-placeholder shapes are cloned at the XML element level.
        - Requires the `python-pptx` library.

    Example:
        combine_ppt(['deck1.pptx', 'deck2.pptx'], 'combined.pptx')
    """

    from pptx import Presentation

    combined_ppt = Presentation()

    for ppt_file in ppt_files:
        current_ppt = Presentation(ppt_file)
        for slide in current_ppt.slides:
            combined_slide = combined_ppt.slides.add_slide(slide.slide_layout)
            for shape in slide.shapes:
                if shape.is_placeholder:
                    new_shape = combined_slide.placeholders[shape.placeholder_format.idx]
                    if shape.has_text_frame:
                        new_shape.text = shape.text
                else:
                    el = shape.element
                    new_el = el.clone()
                    combined_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    combined_ppt.save(output_file)
